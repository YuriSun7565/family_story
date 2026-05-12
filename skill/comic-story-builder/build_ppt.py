"""
Build a kid-friendly storytelling PPT from cropped panels and a captions JSON.

Usage:
    python build_ppt.py <captions.json>
    python build_ppt.py <captions.json> --output story.pptx
    python build_ppt.py <captions.json> --comics-dir path/to/comics

JSON schema:
    {
      "title":    "总封面大标题（中文）",
      "subtitle": "Cover subtitle in English",
      "chapters": [
        {
          "name":     "故事一",
          "subtitle": "Story 1",
          "panels":   ["comic_01.png", "comic_02.png", ...]
        },
        ...
      ],
      "captions": {
        "comic_01.png": {"cn": "中文配文", "en": "English caption"},
        ...
      },
      "theme": {                        // optional override of colours
        "bg":        "#FBF8F1",
        "ink_dark":  "#2C2C2C",
        "ink_brown": "#6B4B2B",
        "accent":    "#C84B31",
        "footer":    "#A08264"
      }
    }
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


THEME_DEFAULT = {
    'bg':        '#FBF8F1',
    'ink_dark':  '#2C2C2C',
    'ink_brown': '#6B4B2B',
    'accent':    '#C84B31',
    'footer':    '#A08264',
}


def _hex(s: str) -> RGBColor:
    s = s.lstrip('#')
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def build_ppt(config: dict, comics_dir: Path, output_path: Path):
    theme = {**THEME_DEFAULT, **(config.get('theme') or {})}
    BG       = _hex(theme['bg'])
    INK_DARK = _hex(theme['ink_dark'])
    INK_BR   = _hex(theme['ink_brown'])
    ACCENT   = _hex(theme['accent'])
    FOOTER   = _hex(theme['footer'])

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # --- helpers --- #
    def fill_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        bg.fill.solid(); bg.fill.fore_color.rgb = BG
        bg.line.fill.background(); bg.shadow.inherit = False
        sp = bg._element; sp.getparent().insert(2, sp)

    def add_text(slide, text, l, t, w, h, *, size=18, bold=False,
                 color=INK_DARK, align=PP_ALIGN.CENTER, font='Microsoft YaHei'):
        if not text:
            return
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font

    def add_image_fit(slide, path, *, top, max_h, max_w_in=11):
        with Image.open(path) as im:
            iw, ih = im.size
        max_w = Inches(max_w_in)
        if iw / ih > max_w / max_h:
            w = max_w; h = int(w * ih / iw)
        else:
            h = max_h; w = int(h * iw / ih)
        slide.shapes.add_picture(str(path), (SW - w) // 2, top, w, h)

    def add_h_stripe(slide, top, height):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SW, height)
        s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

    def add_v_stripe(slide, left, width):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, 0, width, SH)
        s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

    # --- cover slide --- #
    title    = config.get('title', '')
    subtitle = config.get('subtitle', '')
    chapters = config.get('chapters', [])
    captions = config.get('captions', {}) or {}

    if title or subtitle:
        cover = prs.slides.add_slide(blank)
        fill_bg(cover)
        add_h_stripe(cover, Inches(0), Inches(0.3))
        add_h_stripe(cover, Inches(7.2), Inches(0.3))
        if title:
            add_text(cover, title, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.4),
                     size=72, bold=True, color=INK_BR)
        if subtitle:
            add_text(cover, subtitle, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.8),
                     size=32, color=INK_DARK, font='Georgia')
        y = Inches(4.6)
        for ch in chapters:
            line = ch.get('name', '')
            if ch.get('subtitle'):
                line = f'{line}    {ch["subtitle"]}' if line else ch['subtitle']
            if line:
                add_text(cover, '• ' + line, Inches(2), y, Inches(9.3), Inches(0.6),
                         size=20, color=INK_DARK)
                y += Inches(0.55)

    # --- chapters --- #
    n_ch = len(chapters)
    for ci, ch in enumerate(chapters, 1):
        cn = ch.get('name', f'Story {ci}')
        en = ch.get('subtitle', '')
        panels = ch.get('panels', [])

        # chapter cover (skip if there's only one chapter and no name)
        if n_ch > 1 or cn or en:
            cs = prs.slides.add_slide(blank)
            fill_bg(cs)
            add_v_stripe(cs, Inches(0), Inches(0.4))
            if n_ch > 1:
                add_text(cs, f'故事 {ci} / {n_ch}',
                         Inches(1), Inches(2.2), Inches(11.3), Inches(0.6),
                         size=22, color=INK_BR)
            if cn:
                add_text(cs, cn, Inches(1), Inches(2.9), Inches(11.3), Inches(1.6),
                         size=80, bold=True, color=INK_BR)
            if en:
                add_text(cs, en, Inches(1), Inches(4.9), Inches(11.3), Inches(0.8),
                         size=32, color=INK_DARK, font='Georgia')

        for pi, pf in enumerate(panels, 1):
            slide = prs.slides.add_slide(blank)
            fill_bg(slide)
            img_path = comics_dir / pf
            if not img_path.exists():
                print(f'WARNING: missing {img_path}', file=sys.stderr)
                continue
            add_image_fit(slide, img_path, top=Inches(0.5), max_h=Inches(4.6))
            cap = captions.get(pf) or {}
            add_text(slide, cap.get('cn', ''),
                     Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.8),
                     size=28, bold=True, color=INK_BR)
            add_text(slide, cap.get('en', ''),
                     Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7),
                     size=18, color=INK_DARK, font='Georgia')
            if cn:
                add_text(slide, f'{cn}    {pi} / {len(panels)}',
                         Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
                         size=11, color=FOOTER)

    prs.save(str(output_path))
    print(f'Saved: {output_path}')
    print(f'Slides: {len(prs.slides)}')


def main(argv=None):
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument('config', type=Path, help='Path to captions JSON')
    parser.add_argument('--output', '-o', type=Path,
                        help='Output PPTX path (default: <title>.pptx alongside config)')
    parser.add_argument('--comics-dir', type=Path,
                        help='Comics directory (default: <config_dir>/comics)')
    args = parser.parse_args(argv)

    config = json.loads(args.config.read_text(encoding='utf-8'))
    comics_dir = args.comics_dir or (args.config.parent / 'comics')
    if args.output:
        output = args.output
    else:
        name = (config.get('title') or args.config.stem).strip() or 'storybook'
        output = args.config.parent / f'{name}.pptx'

    build_ppt(config, comics_dir, output)
    return 0


if __name__ == '__main__':
    sys.exit(main())
